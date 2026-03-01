from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(22)


def main() -> None:
    prs = Presentation()

    # Slide 1
    add_title_slide(
        prs,
        "MediStock – Phase-wise Brief",
        "Group 57 | Supervisor: Vidyapati Kumar | Ameya Dhurde (2023EBCS392), Karri Leena (2023EBCS343)",
    )

    # Phase 1 (Slides 2-4) - Original proposal
    add_bullets_slide(
        prs,
        "Phase 1: Original Proposal – Problem & Motivation",
        [
            "Hospitals face manual, error-prone medicine inventory tracking",
            "Stockouts impact patient care; overstock increases waste/cost",
            "Need real-time visibility + data-driven procurement planning",
            "Goal: a web-based system for inventory + consumption + alerts",
        ],
    )

    add_bullets_slide(
        prs,
        "Phase 1: Original Proposal – Objectives",
        [
            "Drug catalog management (CRUD) with thresholds",
            "Stock tracking and controlled adjustments",
            "Consumption logging with automatic stock deduction",
            "Alerts + reports for decision support",
            "Forecasting + stockout prediction (AI/ML) for proactive planning",
        ],
    )

    add_bullets_slide(
        prs,
        "Phase 1: Original Proposal – Planned Stack",
        [
            "Backend: Symfony (PHP) using MVC + services",
            "Database: MySQL (Docker for dev)",
            "Frontend: Twig + Bootstrap (responsive UI)",
            "Analytics/Charts: Chart.js",
            "ML: statistical forecasting (moving average + trend), extensible to advanced models",
        ],
    )

    # Phase 2 (Slides 5-7) - Current stage / PoC
    add_bullets_slide(
        prs,
        "Phase 2: Current Stage – What’s Implemented (PoC)",
        [
            "Complete modules: Drugs, Stock, Consumption, Alerts, Reports, Dashboard",
            "Automatic stock deduction on consumption logging",
            "Low-stock detection using threshold logic",
            "Reports: low-stock report + usage report (drug/date filters)",
            "Predictions: consumption forecast + stockout prediction with urgency labels",
        ],
    )

    add_bullets_slide(
        prs,
        "Phase 2: Validation Evidence (Per Supervisor Feedback)",
        [
            "Prediction validation example: actual vs predicted + error metric (MAE/MAPE)",
            "Stock integrity: negative stock prevented (tested + screenshot evidence)",
            "Performance measurements: typical actions under ~2 seconds locally",
            "Testing results reported as pass/fail (not only plans)",
        ],
    )

    add_bullets_slide(
        prs,
        "Phase 2: Architecture Snapshot",
        [
            "MVC flow: Browser → Twig/Bootstrap → Controllers → Services → Doctrine → MySQL",
            "Entities: Drug, Consumption (Drug 1-to-many Consumption)",
            "Services: ConsumptionPredictionService, LowStockPredictionService",
            "Prediction pages: chart + tables (Chart.js)",
            "Security in PoC: CSRF on forms, ORM protections, Twig escaping",
        ],
    )

    # Phase 3 (Slides 8-10) - Future plans
    add_bullets_slide(
        prs,
        "Phase 3: Next Steps – Product Hardening",
        [
            "Authentication + authorization (Admin/Staff roles)",
            "Audit logging + stronger validation for critical actions",
            "Improved notifications (email/SMS) for urgent stockout risk",
            "Deployment planning (production config, backups, monitoring)",
        ],
    )

    add_bullets_slide(
        prs,
        "Phase 3: ML/Analytics Enhancements",
        [
            "Improve forecasting with more data and better evaluation (rolling validation)",
            "Add seasonal pattern detection / advanced models (e.g., LSTM optional)",
            "Category-level analytics, cost/waste insights, procurement recommendations",
            "Dashboards for trends and KPI monitoring over time",
        ],
    )

    add_bullets_slide(
        prs,
        "Phase 3: IoT Integration Plan (Clarified)",
        [
            "Phase 2 PoC does not integrate real hardware (explicitly stated in report)",
            "Phase 3: integrate barcode/RFID inputs for stock updates",
            "Potential smart shelf / ward device feeds into same stock & consumption modules",
            "Define device → API/data format + reliability and security requirements",
        ],
    )

    prs.save("Group57_Vidyapati_MediStock.pptx")


if __name__ == "__main__":
    main()






